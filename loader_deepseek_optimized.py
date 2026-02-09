import os
import sys
import tempfile
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import AutoModel, AutoTokenizer
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# DeepSeek OCR 모델 캐시
_model_cache = None
_tokenizer_cache = None


def _load_model():
    """DeepSeek OCR 모델 로드 (A100 최적화 - 공식 방식)"""
    global _model_cache, _tokenizer_cache

    if _model_cache is None:
        logger.info("DeepSeek OCR 모델 로딩 중 (A100 최적화 - 공식 방식)...")

        model_name = "deepseek-ai/DeepSeek-OCR"

        # GPU 확인
        if not torch.cuda.is_available():
            raise RuntimeError("❌ GPU를 사용할 수 없습니다. 이 스크립트는 GPU 전용입니다.")

        gpu_name = torch.cuda.get_device_name(0)
        gpu_memory = round(torch.cuda.get_device_properties(0).total_memory / 1024**3, 1)
        logger.info(f"GPU: {gpu_name} ({gpu_memory} GB)")

        try:
            # Tokenizer 로드
            _tokenizer_cache = AutoTokenizer.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # Flash Attention 2 사용 가능 여부 확인
            try:
                import flash_attn
                use_flash_attn = True
                logger.info("✓ Flash Attention 2 감지됨")
            except ImportError:
                use_flash_attn = False
                logger.warning("⚠️  Flash Attention 2 미설치 - 기본 attention 사용")
                logger.warning("    설치 권장: pip install flash-attn")

            # 공식 권장 방식으로 모델 로드
            if use_flash_attn:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    _attn_implementation='flash_attention_2',
                    trust_remote_code=True
                )
                logger.info("✓ Flash Attention 2 활성화")
            else:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    trust_remote_code=True
                )

            # 공식 권장: .eval().cuda().to(torch.bfloat16)
            _model_cache = _model_cache.eval().cuda().to(torch.bfloat16)

            logger.info("✓ 모델 로딩 완료 (bfloat16 + CUDA)")

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")
            logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
            logger.error(f"  huggingface-cli download {model_name}")
            raise

    return _model_cache, _tokenizer_cache


def _extract_text_from_image(image: Image.Image, mode: str = "gundam") -> str:
    """
    DeepSeek OCR로 이미지에서 텍스트 추출 (공식 권장 방식)

    Args:
        image: PIL Image 객체
        mode: OCR 모드
            - "tiny": base_size=512, image_size=512, crop_mode=False (가장 빠름)
            - "small": base_size=640, image_size=640, crop_mode=False (빠름)
            - "base": base_size=1024, image_size=1024, crop_mode=False (균형)
            - "large": base_size=1280, image_size=1280, crop_mode=False (고품질)
            - "gundam": base_size=1024, image_size=640, crop_mode=True (권장)

    Returns:
        추출된 텍스트
    """
    model, tokenizer = _load_model()

    # 모드별 설정
    mode_configs = {
        "tiny": {"base_size": 512, "image_size": 512, "crop_mode": False},
        "small": {"base_size": 640, "image_size": 640, "crop_mode": False},
        "base": {"base_size": 1024, "image_size": 1024, "crop_mode": False},
        "large": {"base_size": 1280, "image_size": 1280, "crop_mode": False},
        "gundam": {"base_size": 1024, "image_size": 640, "crop_mode": True},
    }

    config = mode_configs.get(mode, mode_configs["gundam"])
    logger.info(f"OCR 모드: {mode} (base={config['base_size']}, img={config['image_size']}, crop={config['crop_mode']})")

    # 이미지를 임시 파일로 저장
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image.save(tmp.name)
        image_path = tmp.name

    # 임시 출력 디렉토리 생성 (DeepSeek OCR은 output_path가 필수)
    temp_output_dir = tempfile.mkdtemp()

    try:
        # 공식 권장 프롬프트 (한글/영어 최적화)
        prompt = """<image>
Extract all text from this image. The text may contain Korean (한글) and English. Return only the text content without any description.
이 이미지에서 모든 텍스트를 정확하게 추출하세요. 텍스트만 반환하고 설명은 제외하세요."""

        # DeepSeek OCR 추론 (공식 방식)
        response = model.infer(
            tokenizer,
            prompt=prompt,
            image_file=image_path,
            base_size=config["base_size"],
            image_size=config["image_size"],
            crop_mode=config["crop_mode"],
            output_path=temp_output_dir,  # 필수 파라미터
        )

        # 응답이 문자열인 경우 직접 반환
        if isinstance(response, str):
            return response.strip()

        # 응답이 딕셔너리인 경우 텍스트 추출
        if isinstance(response, dict):
            return response.get('text', response.get('output', str(response))).strip()

        # 기타 타입
        return str(response).strip()

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        import traceback
        traceback.print_exc()
        raise
    finally:
        # 임시 파일 삭제
        if os.path.exists(image_path):
            os.unlink(image_path)
        # 임시 출력 디렉토리 삭제
        import shutil
        if os.path.exists(temp_output_dir):
            shutil.rmtree(temp_output_dir, ignore_errors=True)


def _ocr_image(image_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        text = _extract_text_from_image(image, mode=mode)

        # 파일로 저장 (요청된 경우)
        if save_to_file:
            output_path = Path(image_path).stem + "_ocr.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ OCR 결과 저장: {output_path}")

        return text
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_deepseek(file_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm
    import io

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # 고해상도 DPI (A100용)
        dpi = 200

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환
            pix = page.get_pixmap(alpha=False, dpi=dpi)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image, mode=mode)
            if page_text:
                texts.append(f"=== 페이지 {i+1} ===\n{page_text}")

    finally:
        doc.close()

    full_text = "\n\n".join(texts)

    # 파일로 저장 (요청된 경우)
    if save_to_file:
        output_path = Path(file_path).stem + "_ocr.txt"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        logger.info(f"✓ OCR 결과 저장: {output_path}")

    return full_text


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str, mode: str = "gundam", save_to_file: bool = True) -> str:
    """
    DeepSeek OCR로 문서 텍스트 추출 (A100 최적화 - 공식 방식)

    Args:
        file_path: 문서 경로
        mode: OCR 모드 (tiny/small/base/large/gundam)
            - "tiny": 가장 빠름 (512x512)
            - "small": 빠름 (640x640)
            - "base": 균형 (1024x1024)
            - "large": 고품질 (1280x1280)
            - "gundam": 권장 (1024 base + 640 img + crop) ⭐
        save_to_file: True이면 OCR 결과를 파일로 저장 (기본값: True)

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_deepseek(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix == ".docx":
        text = _load_docx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".xlsx":
        text = _load_xlsx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".pptx":
        text = _load_pptx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        # 모드 선택 (옵션)
        mode = "gundam"  # 기본값
        if len(sys.argv) > 2:
            mode = sys.argv[2]

        text = load_document(sys.argv[1], mode=mode, save_to_file=True)
        print("=" * 60)
        print("OCR 처리 완료!")
        print("=" * 60)
        print(f"추출된 텍스트 길이: {len(text)} 자")
        print("=" * 60)
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
        print("=" * 60)
    else:
        print("사용법:")
        print("  python loader_deepseek_optimized.py <파일경로> [모드]")
        print()
        print("모드:")
        print("  tiny   - 가장 빠름 (512x512)")
        print("  small  - 빠름 (640x640)")
        print("  base   - 균형 (1024x1024)")
        print("  large  - 고품질 (1280x1280)")
        print("  gundam - 권장 (기본값, 1024 base + 640 img + crop)")
        print()
        print("예시:")
        print("  python loader_deepseek_optimized.py document.pdf")
        print("  python loader_deepseek_optimized.py document.pdf large")
