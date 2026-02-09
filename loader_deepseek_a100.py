import os
import sys
import tempfile
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import AutoModel, AutoTokenizer
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# DeepSeek OCR 모델 캐시
_model_cache = None
_tokenizer_cache = None


def _load_model():
    """DeepSeek OCR 모델 로드 (A100 최적화)"""
    global _model_cache, _tokenizer_cache

    if _model_cache is None:
        logger.info("DeepSeek OCR 모델 로딩 중 (A100 최적화)...")

        model_name = "deepseek-ai/DeepSeek-OCR"

        # GPU 확인
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"사용 장치: {device}")

        if torch.cuda.is_available():
            gpu_name = torch.cuda.get_device_name(0)
            gpu_memory = round(torch.cuda.get_device_properties(0).total_memory / 1024**3, 1)
            logger.info(f"GPU: {gpu_name} ({gpu_memory} GB)")

        if device == "cpu":
            logger.warning("⚠️  GPU를 사용할 수 없습니다. CPU 모드로 실행됩니다 (매우 느림).")

        try:
            # Tokenizer 로드
            _tokenizer_cache = AutoTokenizer.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # Flash Attention 2 사용 가능 여부 확인
            try:
                import flash_attn
                use_flash_attn = True
                logger.info("Flash Attention 2 감지됨 - 사용 활성화")
            except ImportError:
                use_flash_attn = False
                logger.warning("⚠️  Flash Attention 2 미설치 - 기본 attention 사용")

            # A100 최적화 설정
            if use_flash_attn:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    trust_remote_code=True,
                    torch_dtype=torch.bfloat16,
                    device_map="auto",
                    attn_implementation="flash_attention_2",
                )
                logger.info("✓ 모델 로딩 완료")
                logger.info("✓ A100 최적화 활성화: bfloat16 + Flash Attention 2")
            else:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    trust_remote_code=True,
                    torch_dtype=torch.bfloat16,
                    device_map="auto",
                )
                logger.info("✓ 모델 로딩 완료")
                logger.info("✓ A100 최적화 활성화: bfloat16 (Flash Attention 2 없음)")

            _model_cache.eval()

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")
            logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
            logger.error(f"  huggingface-cli download {model_name}")
            raise

    return _model_cache, _tokenizer_cache


def _extract_text_from_image(image: Image.Image, mode: str = "gundam") -> str:
    """
    DeepSeek OCR로 이미지에서 텍스트 추출 (A100 최적화)

    Args:
        image: PIL Image 객체
        mode: OCR 모드 (A100 최적화: large/gundam)

    Returns:
        추출된 텍스트
    """
    model, tokenizer = _load_model()

    # A100 최적화 모드 설정
    mode_configs = {
        "tiny": {"base_size": 512, "image_size": 512, "crop_mode": False},
        "small": {"base_size": 640, "image_size": 640, "crop_mode": False},
        "base": {"base_size": 1024, "image_size": 1024, "crop_mode": False},
        "large": {"base_size": 1536, "image_size": 1536, "crop_mode": False},  # A100 최적화
        "gundam": {"base_size": 1536, "image_size": 640, "crop_mode": True},  # A100 최적화
    }

    config = mode_configs.get(mode, mode_configs["gundam"])

    # 이미지를 임시 파일로 저장
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image.save(tmp.name)
        image_path = tmp.name

    # output_path 설정: ./pdfs/{임시파일명}.md
    output_dir = Path("./pdfs")
    output_dir.mkdir(parents=True, exist_ok=True)

    # 임시 파일명에서 확장자 제거하고 .md로 설정
    temp_filename = Path(image_path).stem
    output_path = str(output_dir / f"{temp_filename}.md")

    try:
        # 공식 권장 프롬프트 (한글/영어 최적화)
        prompt = "<image>\nExtract all text from this image. The text may contain Korean (한글) and English."

        # 공식 방식: model.infer() 호출
        res = model.infer(
            tokenizer,
            prompt=prompt,
            image_file=image_path,
            output_path=output_path,
            base_size=config["base_size"],
            image_size=config["image_size"],
            crop_mode=config["crop_mode"],
            save_results=True,
            test_compress=True
        )

        # 결과 처리
        if isinstance(res, str):
            return res.strip()
        elif isinstance(res, dict):
            text = res.get('text', res.get('output', res.get('result', '')))
            if text:
                return str(text).strip()

        # output_path에 저장된 파일 확인
        if os.path.exists(output_path):
            with open(output_path, 'r', encoding='utf-8') as f:
                content = f.read().strip()
                if content:
                    return content

        return str(res).strip() if res else ""

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        import traceback
        traceback.print_exc()
        raise
    finally:
        # 임시 파일 삭제
        if os.path.exists(image_path):
            os.unlink(image_path)


def _ocr_image(image_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        text = _extract_text_from_image(image, mode=mode)

        # 파일로 저장 (요청된 경우)
        if save_to_file:
            output_path = Path(image_path).stem + "_ocr.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ OCR 결과 저장: {output_path}")

        return text
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_deepseek(file_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm
    import io

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # A100: 고해상도 DPI 사용 가능
        dpi = 200  # RTX 3080: 150, A100: 200

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환 (고해상도)
            pix = page.get_pixmap(alpha=False, dpi=dpi)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image, mode=mode)
            if page_text:
                texts.append(f"=== 페이지 {i+1} ===\n{page_text}")

    finally:
        doc.close()

    full_text = "\n\n".join(texts)

    # 파일로 저장 (요청된 경우)
    if save_to_file:
        output_path = Path(file_path).stem + "_ocr.txt"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        logger.info(f"✓ OCR 결과 저장: {output_path}")

    return full_text


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str, mode: str = "gundam", save_to_file: bool = True) -> str:
    """
    DeepSeek OCR로 문서 텍스트 추출 (A100 GPU 가속)

    Args:
        file_path: 문서 경로
        mode: OCR 모드 (tiny/small/base/large/gundam)
            - A100 최적화: large (1536x1536) 또는 gundam (1536 base + 640 img)
        save_to_file: True이면 OCR 결과를 파일로 저장 (기본값: True)

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_deepseek(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix == ".docx":
        text = _load_docx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".xlsx":
        text = _load_xlsx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".pptx":
        text = _load_pptx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        # 모드 선택 (옵션)
        mode = "gundam"  # 기본값 (A100 최적화)
        if len(sys.argv) > 2:
            mode = sys.argv[2]

        text = load_document(sys.argv[1], mode=mode, save_to_file=True)
        print("=" * 60)
        print("OCR 처리 완료!")
        print("=" * 60)
        print(f"추출된 텍스트 길이: {len(text)} 자")
        print("=" * 60)
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
        print("=" * 60)
    else:
        print("사용법:")
        print("  python loader_deepseek_a100.py <파일경로> [모드]")
        print()
        print("모드 (A100 최적화):")
        print("  tiny   - 가장 빠름 (512x512)")
        print("  small  - 빠름 (640x640)")
        print("  base   - 균형 (1024x1024)")
        print("  large  - 고품질 A100 (1536x1536)")
        print("  gundam - A100 권장 (기본값, base=1536, img=640, crop=True)")
        print()
        print("예시:")
        print("  python loader_deepseek_a100.py document.pdf")
        print("  python loader_deepseek_a100.py document.pdf large")
