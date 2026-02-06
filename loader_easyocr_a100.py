import os
import sys
import io
import logging
from pathlib import Path
from PIL import Image
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# EasyOCR 모델 캐시
_reader_cache = None


def _load_model():
    """EasyOCR 모델 로드 (A100 최적화)"""
    global _reader_cache

    if _reader_cache is None:
        logger.info("EasyOCR 모델 로딩 중 (A100 최적화)...")

        import easyocr
        import torch

        # GPU 확인
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"사용 장치: {device}")

        if torch.cuda.is_available():
            gpu_name = torch.cuda.get_device_name(0)
            gpu_memory = round(torch.cuda.get_device_properties(0).total_memory / 1024**3, 1)
            logger.info(f"GPU: {gpu_name} ({gpu_memory} GB)")

        # EasyOCR Reader 생성 (한글, 영어)
        _reader_cache = easyocr.Reader(
            ['ko', 'en'],  # 한글, 영어
            gpu=(device == "cuda"),
            verbose=False
        )
        logger.info("✓ EasyOCR 로딩 완료")

    return _reader_cache


def _extract_text_from_image(image: Image.Image) -> str:
    """
    EasyOCR로 이미지에서 텍스트 추출 (A100 최적화)
    """
    reader = _load_model()

    # A100: 더 큰 이미지 처리 가능
    max_size = 2048  # RTX 3080: 1536, A100: 2048
    if max(image.size) > max_size:
        ratio = max_size / max(image.size)
        new_size = tuple(int(dim * ratio) for dim in image.size)
        image = image.resize(new_size, Image.Resampling.LANCZOS)
        logger.info(f"이미지 리사이징: {new_size}")

    try:
        # PIL Image를 numpy array로 변환
        import numpy as np
        image_array = np.array(image)

        # OCR 실행
        results = reader.readtext(
            image_array,
            detail=0,  # 텍스트만 반환
            paragraph=True  # 단락으로 그룹화
        )

        # 결과 조합
        text = '\n'.join(results)
        return text

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        import traceback
        traceback.print_exc()
        raise


def _ocr_image(image_path: str) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        return _extract_text_from_image(image)
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_easyocr(file_path: str) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # A100: 고해상도 DPI 사용 가능
        dpi = 200  # RTX 3080: 150, A100: 200

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환 (고해상도)
            pix = page.get_pixmap(alpha=False, dpi=dpi)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image)
            if page_text:
                texts.append(page_text)

    finally:
        doc.close()

    return "\n\n".join(texts)


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str) -> str:
    """
    EasyOCR로 문서 텍스트 추출 (A100 GPU 가속)

    Args:
        file_path: 문서 경로

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_easyocr(file_path)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path)
    elif suffix == ".docx":
        return _load_docx(file_path)
    elif suffix == ".xlsx":
        return _load_xlsx(file_path)
    elif suffix == ".pptx":
        return _load_pptx(file_path)
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1])
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
