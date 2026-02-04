import os
import sys
import tempfile
import logging
import platform
from pathlib import Path

# PaddleX가 참조하는 구 langchain 경로 호환
if "langchain.docstore.document" not in sys.modules:
    import langchain_core.documents as _lc_docs
    _docmod = type(sys)("langchain.docstore.document")
    _docmod.Document = _lc_docs.Document
    sys.modules["langchain.docstore.document"] = _docmod
    if "langchain.docstore" not in sys.modules:
        _ds = type(sys)("langchain.docstore")
        sys.modules["langchain.docstore"] = _ds
if "langchain.text_splitter" not in sys.modules:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    _tsmod = type(sys)("langchain.text_splitter")
    _tsmod.RecursiveCharacterTextSplitter = RecursiveCharacterTextSplitter
    sys.modules["langchain.text_splitter"] = _tsmod

import fitz  # PyMuPDF
from paddleocr import PaddleOCR
import paddle

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# PaddleOCR 인스턴스 캐시
_ocr_cache = {}

def _is_wsl():
    """WSL 환경 감지"""
    if platform.system() != "Linux":
        return False
    try:
        with open("/proc/version", "r", encoding="utf-8", errors="ignore") as f:
            return "microsoft" in f.read().lower() or "wsl" in f.read().lower()
    except Exception:
        return bool(os.environ.get("WSL_DISTRO_NAME"))

def _setup_gpu():
    """GPU 설정 및 확인"""
    try:
        # WSL 환경 확인
        is_wsl = _is_wsl()
        if is_wsl:
            logger.info("WSL 환경 감지")
            # WSL CUDA 라이브러리 경로 확인
            wsl_cuda_path = "/usr/lib/wsl/lib"
            if os.path.exists(wsl_cuda_path):
                current_ld_path = os.environ.get("LD_LIBRARY_PATH", "")
                if wsl_cuda_path not in current_ld_path:
                    os.environ["LD_LIBRARY_PATH"] = f"{wsl_cuda_path}:{current_ld_path}"
                    logger.info(f"LD_LIBRARY_PATH에 {wsl_cuda_path} 추가")

        # CUDA 사용 가능 여부 확인
        if not paddle.device.is_compiled_with_cuda():
            logger.error("❌ PaddlePaddle이 CUDA 지원 없이 설치되었습니다.")
            logger.error("다음 명령어로 GPU 버전을 설치하세요:")
            logger.error("pip install paddlepaddle-gpu==2.6.1.post118 -f https://www.paddlepaddle.org.cn/whl/linux/mkl/avx/stable.html")
            return False

        # GPU 개수 확인
        gpu_count = paddle.device.cuda.device_count()
        if gpu_count == 0:
            logger.error("❌ GPU가 감지되지 않습니다.")
            logger.error("nvidia-smi로 GPU 상태를 확인하세요.")
            return False

        # GPU 설정
        paddle.device.set_device('gpu:0')
        logger.info(f"✓ GPU 설정 완료 (GPU 개수: {gpu_count})")

        # GPU 테스트
        try:
            test_tensor = paddle.to_tensor([1.0, 2.0, 3.0])
            result = test_tensor + 1
            logger.info(f"✓ GPU 연산 테스트 성공")
            return True
        except Exception as test_err:
            logger.error(f"❌ GPU 테스트 실패: {test_err}")
            return False

    except Exception as e:
        logger.error(f"❌ GPU 설정 중 오류 발생: {e}")
        return False

def _get_ocr(lang="korean"):
    """
    GPU 최적화 OCR 인스턴스 생성

    lang: 인식할 언어
    - "korean" (기본): PP-OCRv4 한글 모델 → 한글 + 영어 동시 지원
    - "ch": 중국어(간체/번체/병음) + 영어 + 일본어
    - "en", "fr", "de", "japan" 등 다국어 지원
    """
    global _ocr_cache

    if lang not in _ocr_cache:
        # GPU 설정 (최초 1회만)
        if not _ocr_cache:
            gpu_ok = _setup_gpu()
            if not gpu_ok:
                raise RuntimeError("GPU 초기화 실패. GPU 환경을 확인하세요.")

        # WSL에서는 MKLDNN 비활성화
        enable_mkldnn = not _is_wsl()

        logger.info(f"PaddleOCR 초기화 중 (언어: {lang}, GPU 모드)")

        try:
            # GPU 최적화 설정
            _ocr_cache[lang] = PaddleOCR(
                use_angle_cls=True,  # 텍스트 방향 분류 사용
                lang=lang,
                use_gpu=True,  # GPU 사용 명시
                gpu_mem=8000,  # RTX 3080은 10GB, 8GB 할당
                enable_mkldnn=enable_mkldnn,
                det_db_thresh=0.3,  # detection threshold
                det_db_box_thresh=0.5,  # box threshold
                det_db_unclip_ratio=1.6,  # unclip ratio
                use_dilation=False,  # dilation 비활성화 (속도 향상)
                det_limit_side_len=960,  # 이미지 크기 제한 (GPU 메모리 절약)
                det_limit_type='max',  # 최대 크기 제한
            )
            logger.info(f"✓ OCR 초기화 완료 (언어: {lang})")
        except Exception as e:
            logger.error(f"❌ OCR 초기화 실패: {e}")
            raise

    return _ocr_cache[lang]

def _extract_text_from_result(result):
    """PaddleOCR 결과에서 텍스트 추출"""
    if not result or len(result) == 0:
        return ""

    lines = []
    for line in result:
        if line is None:
            continue
        # PaddleOCR 2.x 형식: [[box], (text, confidence)]
        if isinstance(line, list) and len(line) >= 2:
            text_info = line[1]
            if isinstance(text_info, (list, tuple)) and len(text_info) >= 1:
                text = str(text_info[0]).strip()
                if text:
                    lines.append(text)

    return "\n".join(lines)

def _ocr_image(ocr, image_path: str) -> str:
    """
    이미지 OCR 처리 (GPU 최적화)
    """
    from PIL import Image

    try:
        # 이미지 크기 확인
        img = Image.open(image_path)
        width, height = img.size
        max_dim = max(width, height)

        logger.info(f"이미지 크기: {width}x{height}")

        # 너무 큰 이미지는 리사이징 (GPU 메모리 절약)
        if max_dim > 1920:
            scale = 1920 / max_dim
            new_width = int(width * scale)
            new_height = int(height * scale)
            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            logger.info(f"이미지 리사이징: {new_width}x{new_height}")

            # 임시 파일로 저장
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                img.save(tmp.name)
                try:
                    result = ocr.ocr(tmp.name, cls=True)
                    if result and len(result) > 0:
                        return _extract_text_from_result(result[0])
                    return ""
                finally:
                    os.unlink(tmp.name)
        else:
            result = ocr.ocr(image_path, cls=True)
            if result and len(result) > 0:
                return _extract_text_from_result(result[0])
            return ""

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        raise

def _load_pdf_with_paddleocr(file_path: str, lang: str = "korean") -> str:
    """PDF에서 텍스트 추출 (GPU 가속)"""
    ocr = _get_ocr(lang)
    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"PDF 페이지 수: {total_pages}")

        for i in range(total_pages):
            logger.info(f"페이지 {i+1}/{total_pages} 처리 중...")
            page = doc.load_page(i)

            # DPI 150으로 이미지 변환 (품질과 속도 균형)
            pix = page.get_pixmap(alpha=False, dpi=150)

            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                pix.save(f.name)
                try:
                    page_text = _ocr_image(ocr, f.name)
                    if page_text:
                        texts.append(page_text)
                finally:
                    os.unlink(f.name)
    finally:
        doc.close()

    return "\n\n".join(texts)

def _load_image_with_paddleocr(file_path: str, lang: str = "korean") -> str:
    """이미지에서 텍스트 추출 (GPU 가속)"""
    ocr = _get_ocr(lang)
    return _ocr_image(ocr, file_path)

def _load_docx(file_path: str) -> str:
    """DOCX 파일 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def _load_xlsx(file_path: str) -> str:
    """XLSX 파일 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)

def _load_pptx(file_path: str) -> str:
    """PPTX 파일 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)

def load_document(file_path: str, lang: str = "korean") -> str:
    """
    GPU 가속 문서 텍스트 추출

    file_path: 문서 경로
    lang: OCR 언어 (기본: korean)
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    logger.info(f"문서 로딩 시작: {file_path}")
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _load_pdf_with_paddleocr(file_path, lang)
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _load_image_with_paddleocr(file_path, lang)
    if suffix == ".docx":
        return _load_docx(file_path)
    if suffix == ".xlsx":
        return _load_xlsx(file_path)
    if suffix == ".pptx":
        return _load_pptx(file_path)

    # 기본: 이미지로 시도
    try:
        return _load_image_with_paddleocr(file_path, lang)
    except Exception:
        raise ValueError(
            f"지원하지 않는 파일 형식: {suffix}. "
            "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
        )

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1])
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
