import os
import sys
import io
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import Qwen3VLForConditionalGeneration, AutoProcessor
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Qwen3-VL 모델 캐시
_model_cache = None
_processor_cache = None


def _load_model():
    """Qwen3-VL 모델 로드 (A100 최적화)"""
    global _model_cache, _processor_cache

    if _model_cache is None:
        logger.info("Qwen3-VL 모델 로딩 중 (A100 최적화)...")

        # A100 최적화: 8B 기본, 32B까지 가능
        # Qwen3-VL-8B: 빠르고 고품질 (기본)
        # Qwen3-VL-32B: 최고 품질 (A100 40GB 이상 권장)
        model_name = "Qwen/Qwen3-VL-8B-Instruct"

        # GPU 확인
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"사용 장치: {device}")

        if torch.cuda.is_available():
            gpu_name = torch.cuda.get_device_name(0)
            gpu_memory = round(torch.cuda.get_device_properties(0).total_memory / 1024**3, 1)
            logger.info(f"GPU: {gpu_name} ({gpu_memory} GB)")

        if device == "cpu":
            logger.warning("⚠️  GPU를 사용할 수 없습니다. CPU 모드로 실행됩니다 (매우 느림).")

        try:
            # Processor 로드
            _processor_cache = AutoProcessor.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # A100 최적화 설정
            # - bfloat16: A100 최적화 데이터 타입
            # - Flash Attention 2: 자동 활성화 (설치된 경우)
            # - device_map="auto": 자동 GPU 할당
            _model_cache = Qwen3VLForConditionalGeneration.from_pretrained(
                model_name,
                trust_remote_code=True,
                torch_dtype=torch.bfloat16,  # A100 최적화
                device_map="auto",
                attn_implementation="flash_attention_2",  # Flash Attention 2 사용
            )

            _model_cache.eval()
            logger.info("✓ 모델 로딩 완료")
            logger.info("✓ A100 최적화 활성화: bfloat16 + Flash Attention 2")

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")

            # Flash Attention 2 없이 재시도
            if "flash_attention_2" in str(e):
                logger.warning("Flash Attention 2 없이 재시도...")
                _model_cache = Qwen3VLForConditionalGeneration.from_pretrained(
                    model_name,
                    trust_remote_code=True,
                    torch_dtype=torch.bfloat16,
                    device_map="auto",
                )
                _model_cache.eval()
                logger.info("✓ 모델 로딩 완료 (Flash Attention 2 없음)")
            else:
                logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
                logger.error(f"  huggingface-cli download {model_name}")
                raise

    return _model_cache, _processor_cache


def _extract_text_from_image(image: Image.Image) -> str:
    """
    Qwen3-VL로 이미지에서 텍스트 추출 (A100 최적화)
    """
    model, processor = _load_model()
    device = next(model.parameters()).device

    # A100: 더 큰 이미지 처리 가능
    max_size = 2048  # RTX 3080: 1280, A100: 2048
    if max(image.size) > max_size:
        ratio = max_size / max(image.size)
        new_size = tuple(int(dim * ratio) for dim in image.size)
        image = image.resize(new_size, Image.Resampling.LANCZOS)
        logger.info(f"이미지 리사이징: {new_size}")

    try:
        # OCR 프롬프트 구성 (Qwen3-VL 스타일)
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "image", "image": image},
                    {
                        "type": "text",
                        "text": """이 이미지에서 모든 텍스트를 정확하게 추출해주세요. 한글과 영어가 포함되어 있을 수 있습니다.
Extract all text from this image accurately. The text contains Korean and English characters.
텍스트의 레이아웃을 그대로 유지하고, 추출된 텍스트만 반환하세요. 추가 설명 없이 텍스트만 출력해주세요.
Please preserve the text layout and return only the extracted text without any additional explanation."""
                    },
                ],
            }
        ]

        # 입력 준비
        text_prompt = processor.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )

        inputs = processor(
            text=[text_prompt],
            images=[image],
            return_tensors="pt",
        )
        inputs = inputs.to(device)

        # 추론 (A100: 더 많은 토큰 생성 가능)
        with torch.no_grad():
            outputs = model.generate(
                **inputs,
                max_new_tokens=4096,  # RTX 3080: 2048, A100: 4096
                do_sample=False,
            )

        # 디코딩 (입력 토큰 제외)
        generated_ids = [
            output_ids[len(input_ids):]
            for input_ids, output_ids in zip(inputs.input_ids, outputs)
        ]

        response = processor.batch_decode(
            generated_ids,
            skip_special_tokens=True,
            clean_up_tokenization_spaces=False
        )[0]

        return response.strip()

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        raise


def _ocr_image(image_path: str) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        return _extract_text_from_image(image)
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_qwen(file_path: str) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # A100: 고해상도 DPI 사용 가능
        dpi = 200  # RTX 3080: 150, A100: 200

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환 (고해상도)
            pix = page.get_pixmap(alpha=False, dpi=dpi)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image)
            if page_text:
                texts.append(page_text)

    finally:
        doc.close()

    return "\n\n".join(texts)


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str) -> str:
    """
    Qwen3-VL로 문서 텍스트 추출 (A100 GPU 가속)

    Args:
        file_path: 문서 경로

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_qwen(file_path)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path)
    elif suffix == ".docx":
        return _load_docx(file_path)
    elif suffix == ".xlsx":
        return _load_xlsx(file_path)
    elif suffix == ".pptx":
        return _load_pptx(file_path)
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1])
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
