import os
import sys
import tempfile
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import AutoModel, AutoTokenizer
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# DeepSeek OCR 모델 캐시
_model_cache = None
_tokenizer_cache = None


def _load_model():
    """DeepSeek OCR 모델 로드"""
    global _model_cache, _tokenizer_cache

    if _model_cache is None:
        logger.info("DeepSeek OCR 모델 로딩 중...")

        model_name = "deepseek-ai/DeepSeek-OCR"

        # GPU 확인
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"사용 장치: {device}")

        if device == "cpu":
            logger.warning("⚠️  GPU를 사용할 수 없습니다. CPU 모드로 실행됩니다 (매우 느림).")

        try:
            # Tokenizer 로드
            _tokenizer_cache = AutoTokenizer.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # Model 로드
            _model_cache = AutoModel.from_pretrained(
                model_name,
                trust_remote_code=True,
                torch_dtype=torch.bfloat16 if device == "cuda" else torch.float32,
                device_map="auto" if device == "cuda" else None,
            )

            if device == "cpu":
                _model_cache = _model_cache.to(device)

            _model_cache.eval()
            logger.info("✓ 모델 로딩 완료")

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")
            logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
            logger.error(f"  huggingface-cli download {model_name}")
            raise

    return _model_cache, _tokenizer_cache


def _extract_text_from_image(image: Image.Image) -> str:
    """
    DeepSeek OCR로 이미지에서 텍스트 추출
    """
    model, tokenizer = _load_model()
    device = next(model.parameters()).device

    # 이미지 크기 조정 (메모리 관리)
    max_size = 1024
    if max(image.size) > max_size:
        ratio = max_size / max(image.size)
        new_size = tuple(int(dim * ratio) for dim in image.size)
        image = image.resize(new_size, Image.Resampling.LANCZOS)
        logger.info(f"이미지 리사이징: {new_size}")

    # 이미지를 임시 파일로 저장
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image.save(tmp.name)
        image_path = tmp.name

    # 임시 출력 디렉토리 생성
    temp_output_dir = tempfile.mkdtemp()

    try:
        # DeepSeek OCR 프롬프트 (한글/영어 최적화)
        prompt = """<image>
Extract all text from this image. The text may contain Korean (한글) and English. Return only the text.
이 이미지에서 모든 텍스트를 정확하게 추출하세요. 텍스트만 반환하고 설명은 제외하세요."""

        # stdout 캡처 (모델이 콘솔에 직접 출력하는 것을 캡처)
        import contextlib
        import io as io_module

        # DeepSeek OCR의 infer 메서드 사용
        # save_results=True로 설정하여 결과를 파일로 저장

        # stdout을 캡처하여 저장
        stdout_capture = io_module.StringIO()
        stderr_capture = io_module.StringIO()

        with contextlib.redirect_stdout(stdout_capture), \
             contextlib.redirect_stderr(stderr_capture):
            result = model.infer(
                tokenizer,
                prompt=prompt,
                image_file=image_path,
                base_size=1024,  # RTX 3080: 1024
                image_size=640,
                crop_mode=True,
                save_results=True,  # True로 변경하여 파일 저장
                output_path=temp_output_dir
            )

        # stdout에서 캡처된 OCR 결과 가져오기
        captured_output = stdout_capture.getvalue()

        # 결과 우선순위:
        # 1. 캡처된 stdout 출력 (실제 OCR 결과)
        # 2. 저장된 텍스트 파일
        # 3. 반환된 result 객체

        if captured_output.strip():
            # stdout에서 캡처된 텍스트 사용
            response = captured_output.strip()
        else:
            # 파일에서 읽기 시도
            import glob
            txt_files = glob.glob(os.path.join(temp_output_dir, "*.txt"))

            if txt_files:
                with open(txt_files[0], 'r', encoding='utf-8') as f:
                    response = f.read()
            elif isinstance(result, dict):
                response = result.get('text', result.get('output', str(result)))
            elif result:
                response = str(result)
            else:
                response = ""

        return response.strip()

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        import traceback
        traceback.print_exc()
        raise
    finally:
        # 임시 파일 삭제
        if os.path.exists(image_path):
            os.unlink(image_path)
        # 임시 출력 디렉토리 삭제
        import shutil
        if os.path.exists(temp_output_dir):
            shutil.rmtree(temp_output_dir, ignore_errors=True)


def _ocr_image(image_path: str) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        return _extract_text_from_image(image)
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_deepseek(file_path: str) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm
    import io

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환
            pix = page.get_pixmap(alpha=False, dpi=150)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image)
            if page_text:
                texts.append(page_text)

    finally:
        doc.close()

    return "\n\n".join(texts)


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str) -> str:
    """
    DeepSeek OCR로 문서 텍스트 추출 (GPU 가속)

    Args:
        file_path: 문서 경로

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_deepseek(file_path)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path)
    elif suffix == ".docx":
        return _load_docx(file_path)
    elif suffix == ".xlsx":
        return _load_xlsx(file_path)
    elif suffix == ".pptx":
        return _load_pptx(file_path)
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


def summarize_text(text: str) -> str:
    """
    DeepSeek 모델을 사용하여 텍스트 요약 생성

    Args:
        text: 요약할 텍스트

    Returns:
        요약된 텍스트
    """
    model, tokenizer = _load_model()
    device = next(model.parameters()).device

    logger.info("텍스트 요약 생성 중...")

    # 텍스트가 너무 길면 앞부분만 사용 (토큰 제한)
    max_chars = 6000  # RTX 3080: 약 1500 토큰
    if len(text) > max_chars:
        text_to_summarize = text[:max_chars] + "..."
        logger.info(f"텍스트가 길어서 처음 {max_chars}자만 요약합니다.")
    else:
        text_to_summarize = text

    # 요약 프롬프트
    prompt = f"""다음 문서의 내용을 한글로 요약해주세요. 주요 내용과 핵심 포인트를 포함하여 3-5개 문단으로 작성해주세요.

문서 내용:
{text_to_summarize}

요약:"""

    try:
        # stdout 억제
        import contextlib
        import io as io_module

        with contextlib.redirect_stdout(io_module.StringIO()), \
             contextlib.redirect_stderr(io_module.StringIO()):

            # 입력 준비
            inputs = tokenizer(prompt, return_tensors="pt")
            inputs = {k: v.to(device) for k, v in inputs.items()}

            # 요약 생성
            with torch.no_grad():
                outputs = model.generate(
                    **inputs,
                    max_new_tokens=768,  # RTX 3080: 768
                    do_sample=False,
                    pad_token_id=tokenizer.eos_token_id if tokenizer.eos_token_id else tokenizer.pad_token_id,
                )

            # 디코딩
            summary = tokenizer.decode(
                outputs[0][inputs["input_ids"].shape[1]:],
                skip_special_tokens=True
            ).strip()

        logger.info("✓ 요약 생성 완료")
        return summary

    except Exception as e:
        logger.error(f"요약 생성 중 오류: {e}")
        return f"요약 생성 실패: {str(e)}"


if __name__ == "__main__":
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1])
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
