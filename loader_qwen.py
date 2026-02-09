import os
import sys
import io
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import Qwen3VLForConditionalGeneration, AutoProcessor
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Qwen3-VL 모델 캐시
_model_cache = None
_processor_cache = None


def _load_model():
    """Qwen3-VL 모델 로드"""
    global _model_cache, _processor_cache

    if _model_cache is None:
        logger.info("Qwen3-VL 모델 로딩 중...")

        # Qwen3-VL-4B 사용 (빠른 속도, RTX 3080 10GB에 최적)
        # 더 높은 품질이 필요하면: "Qwen/Qwen3-VL-8B-Instruct" (느림)
        model_name = "Qwen/Qwen3-VL-4B-Instruct"

        # GPU 확인
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"사용 장치: {device}")

        if device == "cpu":
            logger.warning("⚠️  GPU를 사용할 수 없습니다. CPU 모드로 실행됩니다 (매우 느림).")

        try:
            # Processor 로드
            _processor_cache = AutoProcessor.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # Model 로드
            _model_cache = Qwen3VLForConditionalGeneration.from_pretrained(
                model_name,
                trust_remote_code=True,
                torch_dtype=torch.bfloat16 if device == "cuda" else torch.float32,
                device_map="auto" if device == "cuda" else None,
            )

            if device == "cpu":
                _model_cache = _model_cache.to(device)

            _model_cache.eval()
            logger.info("✓ 모델 로딩 완료")

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")
            logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
            logger.error(f"  huggingface-cli download {model_name}")
            raise

    return _model_cache, _processor_cache


def _extract_text_from_image(image: Image.Image) -> str:
    """
    Qwen3-VL로 이미지에서 텍스트 추출
    """
    model, processor = _load_model()
    device = next(model.parameters()).device

    # 이미지 크기 조정 (메모리 관리)
    max_size = 1280
    if max(image.size) > max_size:
        ratio = max_size / max(image.size)
        new_size = tuple(int(dim * ratio) for dim in image.size)
        image = image.resize(new_size, Image.Resampling.LANCZOS)
        logger.info(f"이미지 리사이징: {new_size}")

    try:
        # OCR 프롬프트 구성 (Qwen3-VL 스타일)
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "image", "image": image},
                    {
                        "type": "text",
                        "text": """이 이미지에서 모든 텍스트를 정확하게 추출해주세요. 한글과 영어가 포함되어 있을 수 있습니다.
Extract all text from this image accurately. The text contains Korean and English characters.
텍스트의 레이아웃을 그대로 유지하고, 추출된 텍스트만 반환하세요. 추가 설명 없이 텍스트만 출력해주세요.
Please preserve the text layout and return only the extracted text without any additional explanation."""
                    },
                ],
            }
        ]

        # 입력 준비
        text_prompt = processor.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )

        inputs = processor(
            text=[text_prompt],
            images=[image],
            return_tensors="pt",
        )
        inputs = inputs.to(device)

        # 추론
        with torch.no_grad():
            outputs = model.generate(
                **inputs,
                max_new_tokens=2048,
                do_sample=False,
            )

        # 디코딩 (입력 토큰 제외)
        generated_ids = [
            output_ids[len(input_ids):]
            for input_ids, output_ids in zip(inputs.input_ids, outputs)
        ]

        response = processor.batch_decode(
            generated_ids,
            skip_special_tokens=True,
            clean_up_tokenization_spaces=False
        )[0]

        return response.strip()

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        raise


def _ocr_image(image_path: str, save_to_file: bool = False) -> str:
    """이미지 파일 OCR"""
    logger.info(f"이미지 OCR: {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        text = _extract_text_from_image(image)

        # 파일로 저장 (요청된 경우)
        if save_to_file:
            output_path = Path(image_path).stem + "_ocr.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ OCR 결과 저장: {output_path}")

        return text
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_qwen(file_path: str, save_to_file: bool = False) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환
            pix = page.get_pixmap(alpha=False, dpi=150)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image)
            if page_text:
                texts.append(f"=== 페이지 {i+1} ===\n{page_text}")

    finally:
        doc.close()

    full_text = "\n\n".join(texts)

    # 파일로 저장 (요청된 경우)
    if save_to_file:
        output_path = Path(file_path).stem + "_ocr.txt"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        logger.info(f"✓ OCR 결과 저장: {output_path}")

    return full_text


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str, save_to_file: bool = True) -> str:
    """
    Qwen3-VL로 문서 텍스트 추출 (GPU 가속)

    Args:
        file_path: 문서 경로
        save_to_file: True이면 OCR 결과를 파일로 저장 (기본값: True)

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_qwen(file_path, save_to_file=save_to_file)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path, save_to_file=save_to_file)
    elif suffix == ".docx":
        text = _load_docx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".xlsx":
        text = _load_xlsx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".pptx":
        text = _load_pptx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path, save_to_file=save_to_file)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1], save_to_file=True)
        print("=" * 60)
        print("OCR 처리 완료!")
        print("=" * 60)
        print(f"추출된 텍스트 길이: {len(text)} 자")
        print("=" * 60)
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
        print("=" * 60)
