#!/usr/bin/env python3
"""
DeepSeek-OCR Official Implementation
Based on: https://huggingface.co/deepseek-ai/DeepSeek-OCR

Official Dependencies:
- torch==2.6.0
- transformers==4.46.3
- tokenizers==0.20.3
- flash-attn (optional, for acceleration)
"""

import os
import sys
import tempfile
import logging
from pathlib import Path
import torch
from PIL import Image
from transformers import AutoModel, AutoTokenizer
import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# DeepSeek OCR 모델 캐시
_model_cache = None
_tokenizer_cache = None


def _load_model():
    """DeepSeek OCR 모델 로드 (공식 방식)"""
    global _model_cache, _tokenizer_cache

    if _model_cache is None:
        logger.info("DeepSeek OCR 모델 로딩 중 (공식 구현)...")

        model_name = "deepseek-ai/DeepSeek-OCR"

        # GPU 확인
        if not torch.cuda.is_available():
            raise RuntimeError("❌ GPU를 사용할 수 없습니다. DeepSeek-OCR은 GPU 전용입니다.")

        gpu_name = torch.cuda.get_device_name(0)
        gpu_memory = round(torch.cuda.get_device_properties(0).total_memory / 1024**3, 1)
        logger.info(f"GPU: {gpu_name} ({gpu_memory} GB)")

        try:
            # Tokenizer 로드 (공식 방식)
            _tokenizer_cache = AutoTokenizer.from_pretrained(
                model_name,
                trust_remote_code=True
            )

            # Flash Attention 2 사용 가능 여부 확인
            try:
                import flash_attn
                use_flash_attn = True
                logger.info("✓ Flash Attention 2 감지됨")
            except ImportError:
                use_flash_attn = False
                logger.warning("⚠️  Flash Attention 2 미설치 - 기본 attention 사용")

            # 공식 권장 방식으로 모델 로드
            if use_flash_attn:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    _attn_implementation='flash_attention_2',
                    trust_remote_code=True,
                    use_safetensors=True
                )
                logger.info("✓ Flash Attention 2 활성화")
            else:
                _model_cache = AutoModel.from_pretrained(
                    model_name,
                    trust_remote_code=True,
                    use_safetensors=True
                )

            # 공식 예제: model.eval().cuda().to(torch.bfloat16)
            _model_cache = _model_cache.eval().cuda().to(torch.bfloat16)

            logger.info("✓ 모델 로딩 완료 (bfloat16 + CUDA)")

        except Exception as e:
            logger.error(f"❌ 모델 로딩 실패: {e}")
            logger.error("다음 명령어로 수동 다운로드를 시도하세요:")
            logger.error(f"  huggingface-cli download {model_name}")
            raise

    return _model_cache, _tokenizer_cache


def _extract_text_from_image(image: Image.Image, mode: str = "gundam", task: str = "ocr") -> str:
    """
    DeepSeek OCR로 이미지에서 텍스트 추출 (공식 방식 그대로 사용)

    Args:
        image: PIL Image 객체
        mode: OCR 모드
            - "tiny": base_size=512, image_size=512
            - "small": base_size=640, image_size=640
            - "base": base_size=1024, image_size=1024
            - "large": base_size=1280, image_size=1280
            - "gundam": base_size=1024, image_size=640, crop_mode=True (공식 예제)
        task: 작업 유형
            - "ocr": 텍스트 추출
            - "markdown": 마크다운 변환

    Returns:
        추출된 텍스트
    """
    model, tokenizer = _load_model()

    # 모드별 설정 (공식 문서 기준)
    mode_configs = {
        "tiny": {"base_size": 512, "image_size": 512, "crop_mode": False},
        "small": {"base_size": 640, "image_size": 640, "crop_mode": False},
        "base": {"base_size": 1024, "image_size": 1024, "crop_mode": False},
        "large": {"base_size": 1280, "image_size": 1280, "crop_mode": False},
        "gundam": {"base_size": 1024, "image_size": 640, "crop_mode": True},  # 공식 예제
    }

    config = mode_configs.get(mode, mode_configs["gundam"])

    # 이미지를 임시 파일로 저장
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image.save(tmp.name)
        image_path = tmp.name

    # 임시 출력 디렉토리 생성
    temp_output_dir = tempfile.mkdtemp()

    try:
        # 공식 예제 프롬프트 형식
        if task == "markdown":
            # 공식 예제: "<image>\n<|grounding|>Convert the document to markdown. "
            prompt = "<image>\n<|grounding|>Convert the document to markdown."
        else:
            # OCR 전용 (한글/영어 최적화)
            prompt = "<image>\nExtract all text from this image. The text may contain Korean (한글) and English."

        # 공식 예제 그대로: model.infer() 호출
        # res = model.infer(tokenizer, prompt=prompt, image_file=image_file,
        #                   output_path=output_path, base_size=1024, image_size=640,
        #                   crop_mode=True, save_results=True, test_compress=True)

        res = model.infer(
            tokenizer,
            prompt=prompt,
            image_file=image_path,
            output_path=temp_output_dir,
            base_size=config["base_size"],
            image_size=config["image_size"],
            crop_mode=config["crop_mode"],
            save_results=True,  # 공식 예제 참고
            test_compress=True  # 공식 예제 참고
        )

        # 결과 처리
        if isinstance(res, str):
            return res.strip()
        elif isinstance(res, dict):
            # 딕셔너리에서 텍스트 추출
            text = res.get('text', res.get('output', res.get('result', '')))
            if text:
                return str(text).strip()

        # output_path에 저장된 파일 확인
        if os.path.exists(temp_output_dir):
            result_files = list(Path(temp_output_dir).glob('*.txt'))
            if result_files:
                # 첫 번째 결과 파일 읽기
                with open(result_files[0], 'r', encoding='utf-8') as f:
                    content = f.read().strip()
                    if content:
                        return content

        return str(res).strip() if res else ""

    except Exception as e:
        logger.error(f"OCR 처리 중 오류: {e}")
        import traceback
        traceback.print_exc()
        raise
    finally:
        # 임시 파일 삭제
        if os.path.exists(image_path):
            os.unlink(image_path)
        # 임시 출력 디렉토리 삭제
        import shutil
        if os.path.exists(temp_output_dir):
            shutil.rmtree(temp_output_dir, ignore_errors=True)


def _ocr_image(image_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """이미지 파일 OCR"""
    logger.info(f"((((((이미지 OCR)))))): {image_path}")

    try:
        image = Image.open(image_path).convert("RGB")
        text = _extract_text_from_image(image, mode=mode)
        logger.info(f"**********text: {text}")

        # 파일로 저장 (요청된 경우)
        if save_to_file:
            output_path = Path(image_path).stem + "_ocr.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ OCR 결과 저장: {output_path}")

        return text
    except Exception as e:
        logger.error(f"이미지 로드 실패: {e}")
        raise


def _load_pdf_with_deepseek(file_path: str, mode: str = "gundam", save_to_file: bool = False) -> str:
    """PDF 파일 OCR"""
    from tqdm import tqdm
    import io

    logger.info(f"PDF 파일 OCR: {file_path}")

    doc = fitz.open(file_path)
    texts = []

    try:
        total_pages = len(doc)
        logger.info(f"총 {total_pages} 페이지")

        # 고해상도 DPI
        dpi = 200

        # tqdm 프로그레스바 사용
        for i in tqdm(range(total_pages), desc="PDF 페이지 처리", unit="page"):
            page = doc.load_page(i)

            # 페이지를 이미지로 변환
            pix = page.get_pixmap(alpha=False, dpi=dpi)

            # PIL Image로 변환
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))

            # OCR
            page_text = _extract_text_from_image(image, mode=mode)
            if page_text:
                texts.append(f"=== 페이지 {i+1} ===\n{page_text}")

    finally:
        doc.close()

    full_text = "\n\n".join(texts)

    # 파일로 저장 (요청된 경우)
    if save_to_file:
        output_path = Path(file_path).stem + "_ocr.txt"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        logger.info(f"✓ OCR 결과 저장: {output_path}")

    return full_text


def _load_docx(file_path: str) -> str:
    """DOCX 텍스트 추출"""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(file_path: str) -> str:
    """XLSX 텍스트 추출"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    """PPTX 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def load_document(file_path: str, mode: str = "gundam", save_to_file: bool = True) -> str:
    """
    DeepSeek OCR로 문서 텍스트 추출 (공식 구현)

    Args:
        file_path: 문서 경로
        mode: OCR 모드 (tiny/small/base/large/gundam)
        save_to_file: True이면 OCR 결과를 파일로 저장 (기본값: True)

    Returns:
        추출된 텍스트
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    path = Path(file_path)
    suffix = path.suffix.lower()

    logger.info(f"문서 로딩: {file_path}")

    if suffix == ".pdf":
        return _load_pdf_with_deepseek(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
    elif suffix == ".docx":
        text = _load_docx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".xlsx":
        text = _load_xlsx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    elif suffix == ".pptx":
        text = _load_pptx(file_path)
        if save_to_file:
            output_path = path.stem + "_text.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"✓ 텍스트 저장: {output_path}")
        return text
    else:
        # 기본: 이미지로 시도
        try:
            return _ocr_image(file_path, mode=mode, save_to_file=save_to_file)
        except:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix}. "
                "지원 형식: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
            )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        # 모드 선택 (옵션)
        mode = "gundam"  # 기본값 (공식 예제)
        if len(sys.argv) > 2:
            mode = sys.argv[2]

        text = load_document(sys.argv[1], mode=mode, save_to_file=True)
        print("=" * 60)
        print("OCR 처리 완료!")
        print("=" * 60)
        print(f"추출된 텍스트 길이: {len(text)} 자")
        print("=" * 60)
        print("--- 추출된 텍스트 (처음 500자) ---")
        print(text[:500])
        print("=" * 60)
    else:
        print("사용법:")
        print("  python loader_deepseek_official.py <파일경로> [모드]")
        print()
        print("모드 (공식 문서 기준):")
        print("  tiny   - 가장 빠름 (512x512)")
        print("  small  - 빠름 (640x640)")
        print("  base   - 균형 (1024x1024)")
        print("  large  - 고품질 (1280x1280)")
        print("  gundam - 공식 예제 (기본값, base=1024, img=640, crop=True)")
        print()
        print("예시:")
        print("  python loader_deepseek_official.py document.pdf")
        print("  python loader_deepseek_official.py document.pdf large")
