import os
import sys
import tempfile
import logging
import platform
from pathlib import Path

# PaddleX가 참조하는 구 langchain 경로 호환
if "langchain.docstore.document" not in sys.modules:
    import langchain_core.documents as _lc_docs
    _docmod = type(sys)("langchain.docstore.document")
    _docmod.Document = _lc_docs.Document
    sys.modules["langchain.docstore.document"] = _docmod
    if "langchain.docstore" not in sys.modules:
        _ds = type(sys)("langchain.docstore")
        sys.modules["langchain.docstore"] = _ds
if "langchain.text_splitter" not in sys.modules:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    _tsmod = type(sys)("langchain.text_splitter")
    _tsmod.RecursiveCharacterTextSplitter = RecursiveCharacterTextSplitter
    sys.modules["langchain.text_splitter"] = _tsmod

import fitz  # PyMuPDF
from paddleocr import PaddleOCR

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# PaddleOCR 인스턴스 (PaddleOCR 3.x / PaddleX API)
# lang별 캐시: 같은 lang이면 재사용 (다국어 문서 대비)
_ocr_cache = {}

def _is_wsl():
    """WSL(Windows Subsystem for Linux) 여부. WSL에서도 oneDNN 에러 발생."""
    if platform.system() != "Linux":
        return False
    try:
        with open("/proc/version", "r", encoding="utf-8", errors="ignore") as f:
            return "microsoft" in f.read().lower()
    except Exception:
        return bool(os.environ.get("WSL_DISTRO_NAME"))

def _get_ocr(lang="korean"):
    """
    lang: 인식할 언어. 하나만 지정 가능하지만, 일부 모델은 여러 언어를 함께 지원합니다.
    - "korean" (기본): PP-OCRv5 한글 모델 → 한글 + 영어 동시 지원
    - "ch": 중국어(간체/번체/병음) + 영어 + 일본어
    - "en", "fr", "de", "japan" 등 106개 언어 약어 지원 (문서 Section 4 참고)
    """
    global _ocr_cache
    if lang not in _ocr_cache:
        # GPU 사용 가능 여부 확인 및 device 설정
        # 환경 변수로 GPU 사용 강제 비활성화 가능: DOC_OCR_FORCE_CPU=1
        use_gpu_device = False
        force_cpu = os.environ.get("DOC_OCR_FORCE_CPU", "").strip() in ("1", "true", "yes")

        if force_cpu:
            logger.info("GPU disabled by DOC_OCR_FORCE_CPU environment variable")
        else:
            try:
                import paddle
                # CUDA 사용 가능 여부 확인
                if paddle.device.is_compiled_with_cuda():
                    gpu_count = paddle.device.cuda.device_count()
                    if gpu_count > 0:
                        try:
                            paddle.device.set_device('gpu:0')
                            # 간단한 연산으로 GPU 실제 작동 확인
                            test_tensor = paddle.to_tensor([1.0])
                            _ = test_tensor + 1
                            use_gpu_device = True
                            logger.info(f"GPU detected: Using CUDA (device count: {gpu_count})")
                        except Exception as gpu_err:
                            logger.warning(f"GPU available but test failed, fallback to CPU: {gpu_err}")
                            paddle.device.set_device('cpu')
                    else:
                        logger.info("CUDA compiled but no GPU found: Using CPU")
                        paddle.device.set_device('cpu')
                else:
                    logger.info("PaddlePaddle not compiled with CUDA: Using CPU")
                    paddle.device.set_device('cpu')
            except Exception as e:
                logger.info(f"Unable to detect GPU, using CPU: {e}")
                try:
                    import paddle
                    paddle.device.set_device('cpu')
                except:
                    pass

        # Windows·WSL: oneDNN 경로에서 ConvertPirAttribute2RuntimeAttribute 미구현 에러 방지
        # DOC_OCR_DISABLE_MKLDNN=1 이면 어떤 환경에서도 MKLDNN 비활성화
        if os.environ.get("DOC_OCR_DISABLE_MKLDNN", "").strip() in ("1", "true", "yes"):
            _enable_mkldnn = False
        else:
            _enable_mkldnn = platform.system() == "Darwin" or (
                platform.system() == "Linux" and not _is_wsl()
            )

        # 이미지 최대 크기 제한 (기본 960 → 2000으로 증가하여 품질 향상)
        # 큰 이미지는 자동으로 리사이징되어 메모리 사용량 관리
        # PaddleOCR 버전에 따라 파라미터명이 다름: det_limit_side_len (v3.x), max_side_len (구버전)
        ocr_params = {
            "use_doc_orientation_classify": False,
            "use_doc_unwarping": False,
            "use_textline_orientation": False,
            "lang": lang,
            "ocr_version": "PP-OCRv5",  # v5: 한글 모델이 한·영 동시 지원, 106개 언어
            "enable_mkldnn": _enable_mkldnn,
            "det_limit_side_len": 2000,  # PaddleOCR 3.x: detection 이미지 크기 제한
        }

        _ocr_cache[lang] = PaddleOCR(**ocr_params)
    return _ocr_cache[lang]

def _extract_text_from_paddle3_result(result_list):
    """PaddleOCR 3.x predict() 결과에서 텍스트만 추출 (rec_texts)."""
    if not result_list:
        return ""
    lines = []
    for res in result_list:
        data = getattr(res, "json", None) or (res if isinstance(res, dict) else {})
        res_inner = (data.get("res") or data) if isinstance(data, dict) else {}
        rec_texts = res_inner.get("rec_texts") or []
        for t in rec_texts:
            if t and str(t).strip():
                lines.append(str(t).strip())
    return "\n".join(lines)

def _ocr_image(ocr, image_path: str) -> str:
    """
    이미지를 OCR 처리. 너무 큰 이미지는 미리 리사이징하여 메모리 오류 방지.
    """
    from PIL import Image

    # 이미지 크기 확인 및 필요시 리사이징
    try:
        img = Image.open(image_path)
        width, height = img.size
        max_dimension = max(width, height)

        # 2000px 초과 시 리사이징 (메모리 절약)
        if max_dimension > 2000:
            scale = 2000 / max_dimension
            new_width = int(width * scale)
            new_height = int(height * scale)
            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            # 임시 파일로 저장
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                img.save(tmp.name)
                try:
                    result = ocr.predict(tmp.name)
                    return _extract_text_from_paddle3_result(result)
                finally:
                    os.unlink(tmp.name)
        else:
            result = ocr.predict(image_path)
            return _extract_text_from_paddle3_result(result)
    except Exception as e:
        logger.warning(f"이미지 전처리 실패, 원본으로 시도: {e}")
        result = ocr.predict(image_path)
        return _extract_text_from_paddle3_result(result)

def _load_pdf_with_paddleocr(file_path: str, lang: str = "korean") -> str:
    ocr = _get_ocr(lang)
    doc = fitz.open(file_path)
    texts = []
    try:
        for i in range(len(doc)):
            page = doc.load_page(i)
            # DPI 150 → 120으로 낮춤 (메모리 절약, 여전히 충분한 품질)
            pix = page.get_pixmap(alpha=False, dpi=120)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                pix.save(f.name)
                try:
                    page_text = _ocr_image(ocr, f.name)
                    if page_text:
                        texts.append(page_text)
                finally:
                    os.unlink(f.name)
    finally:
        doc.close()
    return "\n\n".join(texts)

def _load_image_with_paddleocr(file_path: str, lang: str = "korean") -> str:
    ocr = _get_ocr(lang)
    return _ocr_image(ocr, file_path)

def _load_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def _load_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)

def _load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)

def load_document(file_path: str, lang: str = "korean") -> str:
    """
    PaddleOCR(이미지/PDF) 또는 기본 라이브러리(DOCX, XLSX, PPTX)로 문서 텍스트 추출.

    lang: OCR 언어. 기본 "korean"은 한글+영어 동시 지원(PP-OCRv5).
          "ch"(중국어+영어+일본어), "en", "fr", "de" 등 106개 언어 약어 지원.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    logger.info(f"Loading document: {file_path}")
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _load_pdf_with_paddleocr(file_path, lang)
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
        return _load_image_with_paddleocr(file_path, lang)
    if suffix == ".docx":
        return _load_docx(file_path)
    if suffix == ".xlsx":
        return _load_xlsx(file_path)
    if suffix == ".pptx":
        return _load_pptx(file_path)

    # 기본: 이미지로 시도 후 실패 시 에러
    try:
        return _load_image_with_paddleocr(file_path, lang)
    except Exception:
        raise ValueError(
            f"Unsupported file type: {suffix}. "
            "Supported: PDF, PNG, JPG, BMP, GIF, WEBP, TIFF, DOCX, XLSX, PPTX."
        )

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        text = load_document(sys.argv[1])
        print("--- Extracted Text (First 500 chars) ---")
        print(text[:500])
